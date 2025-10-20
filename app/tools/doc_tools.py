from pathlib import Path
import os
from typing import Annotated, Dict, List, Optional

import docx
import openpyxl
from langchain_core.tools import tool  # type: ignore
from pptx import Presentation
from typing_extensions import TypedDict

# Define a permanent directory for output files
WORKING_DIRECTORY = Path.cwd() / "media"
# Ensure the directory exists
os.makedirs(WORKING_DIRECTORY, exist_ok=True)


@tool
def create_outline(
    points: Annotated[List[str], "List of main points or sections."],
    file_name: Annotated[str, "File path to save the outline."],
) -> Annotated[str, "Path of the saved outline file."]:
    """Create and save an outline."""
    with (WORKING_DIRECTORY / file_name).open("w") as file:
        for i, point in enumerate(points):
            file.write(f"{i + 1}. {point}\n")
    return f"Outline saved to {file_name}"


@tool
def read_document(
    file_name: Annotated[str, "File path to read the document from."],
    start: Annotated[Optional[int], "The start line. Default is 0"] = None,
    end: Annotated[Optional[int], "The end line. Default is None"] = None,
) -> str:
    """Read the specified document."""
    with (WORKING_DIRECTORY / file_name).open("r") as file:
        lines = file.readlines()
    if start is None:
        start = 0
    return "\n".join(lines[start:end])


@tool
def write_document(
    content: Annotated[str, "Text content to be written into the document."],
    file_name: Annotated[str, "File path to save the document."],
) -> Annotated[str, "Path of the saved document file."]:
    """Create and save a text document."""
    with (WORKING_DIRECTORY / file_name).open("w") as file:
        file.write(content)
    return f"Document saved to {file_name}"


@tool
def edit_document(
    file_name: Annotated[str, "Path of the document to be edited."],
    inserts: Annotated[
        Dict[int, str],
        "Dictionary where key is the line number (1-indexed) and value is the text to be inserted at that line.",
    ],
) -> Annotated[str, "Path of the edited document file."]:
    """Edit a document by inserting text at specific line numbers."""

    with (WORKING_DIRECTORY / file_name).open("r") as file:
        lines = file.readlines()

    sorted_inserts = sorted(inserts.items())

    for line_number, text in sorted_inserts:
        if 1 <= line_number <= len(lines) + 1:
            lines.insert(line_number - 1, text + "\n")
        else:
            return f"Error: Line number {line_number} is out of range."

    with (WORKING_DIRECTORY / file_name).open("w") as file:
        file.writelines(lines)

    return f"Document edited and saved to {file_name}"


@tool
def create_word_document(
    file_name: Annotated[str, "File path to save the document. Should end with .docx"],
    content: Annotated[str, "Text content to be written into the document."],
) -> Annotated[str, "Path of the saved document file."]:
    """Create and save a Word document."""
    if not file_name.endswith(".docx"):
        return "Error: File name must end with .docx"
    doc = docx.Document()
    doc.add_paragraph(content)
    doc.save(WORKING_DIRECTORY / file_name) # type: ignore
    return f"Word document saved to {file_name}"


@tool
def create_excel_spreadsheet(
    file_name: Annotated[
        str, "File path to save the spreadsheet. Should end with .xlsx"
    ],
    data: Annotated[
        List[List[str]], "A list of lists representing rows and columns."
    ],
) -> Annotated[str, "Path of the saved spreadsheet file."]:
    """Create and save an Excel spreadsheet."""
    if not file_name.endswith(".xlsx"):
        return "Error: File name must end with .xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    for row_data in data:
        sheet.append(row_data) # type: ignore
    workbook.save(WORKING_DIRECTORY / file_name)
    return f"Excel spreadsheet saved to {file_name}"


@tool
def create_powerpoint_presentation(
    file_name: Annotated[
        str, "File path to save the presentation. Should end with .pptx"
    ],
    slides_data: Annotated[
        List[Dict[str, str]],
        "A list of dictionaries, where each dictionary has a 'title' and 'content' key.",
    ],
) -> Annotated[str, "Path of the saved presentation file."]:
    """Create and save a PowerPoint presentation."""
    if not file_name.endswith(".pptx"):
        return "Error: File name must end with .pptx"
    prs = Presentation()
    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = slide_info.get("title", "") # type: ignore
        body.text = slide_info.get("content", "") # type: ignore
    prs.save(WORKING_DIRECTORY / file_name) # type: ignore
    return f"PowerPoint presentation saved to {file_name}"
